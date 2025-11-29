"""
File processing utilities for extracting text from various file formats.
"""

import io
from typing import Dict, Any
from pathlib import Path

# PDF processing
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Word document processing
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# Excel processing
try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

# PowerPoint processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class FileProcessor:
    """Process various file types and extract text content."""

    SUPPORTED_EXTENSIONS = {
        '.txt': 'text',
        '.md': 'text',
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'docx',
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.pptx': 'pptx',
        '.ppt': 'pptx',
    }

    @classmethod
    def extract_text(cls, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract text from a file.

        Args:
            file_content: Raw file bytes
            filename: Original filename

        Returns:
            Dictionary with extracted text and metadata
        """
        file_ext = Path(filename).suffix.lower()
        
        if file_ext not in cls.SUPPORTED_EXTENSIONS:
            return {
                "success": False,
                "error": f"Unsupported file type: {file_ext}",
                "supported_types": list(cls.SUPPORTED_EXTENSIONS.keys())
            }

        file_type = cls.SUPPORTED_EXTENSIONS[file_ext]

        try:
            if file_type == 'text':
                text = cls._extract_text_plain(file_content)
            elif file_type == 'pdf':
                text = cls._extract_text_pdf(file_content)
            elif file_type == 'docx':
                text = cls._extract_text_docx(file_content)
            elif file_type == 'excel':
                text = cls._extract_text_excel(file_content)
            elif file_type == 'pptx':
                text = cls._extract_text_pptx(file_content)
            else:
                return {
                    "success": False,
                    "error": f"Handler not implemented for {file_type}"
                }

            return {
                "success": True,
                "text": text,
                "filename": filename,
                "file_type": file_type,
                "char_count": len(text)
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Error processing {filename}: {str(e)}"
            }

    @staticmethod
    def _extract_text_plain(content: bytes) -> str:
        """Extract text from plain text files."""
        # Try different encodings
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode text file")

    @staticmethod
    def _extract_text_pdf(content: bytes) -> str:
        """Extract text from PDF files."""
        if not PDF_AVAILABLE:
            raise ImportError("PyPDF2 not installed. Install with: pip install PyPDF2")

        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)
        
        text_parts = []
        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            if text.strip():
                text_parts.append(f"[Page {page_num + 1}]\n{text}")

        return "\n\n".join(text_parts)

    @staticmethod
    def _extract_text_docx(content: bytes) -> str:
        """Extract text from Word documents."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")

        doc_file = io.BytesIO(content)
        doc = Document(doc_file)
        
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        return "\n\n".join(text_parts)

    @staticmethod
    def _extract_text_excel(content: bytes) -> str:
        """Extract text from Excel files."""
        if not EXCEL_AVAILABLE:
            raise ImportError("openpyxl not installed. Install with: pip install openpyxl")

        excel_file = io.BytesIO(content)
        workbook = load_workbook(excel_file, data_only=True)
        
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"[Sheet: {sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text_parts.append(row_text)

        return "\n\n".join(text_parts)

    @staticmethod
    def _extract_text_pptx(content: bytes) -> str:
        """Extract text from PowerPoint files."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        pptx_file = io.BytesIO(content)
        presentation = Presentation(pptx_file)
        
        text_parts = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_parts.append(f"[Slide {slide_num}]")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

        return "\n\n".join(text_parts)


# Convenience instance
file_processor = FileProcessor()
